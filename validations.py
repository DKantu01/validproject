"""
The flask application package.
Pathnames are updated via if statements in the code. 
We use postgreSQL instead of MS SQL. 
"""

from flask import Flask, render_template_string, render_template, redirect, url_for, current_app, request, session
from cProfile import label
from html.entities import html5
from textwrap import wrap
from tkinter.ttk import Style
from turtle import width
from xml.etree.ElementTree import tostring
from xml.sax.handler import feature_validation
from optparse import Option, check_builtin
from pickle import FALSE
import dash
from dash import Dash, dash_table, dcc, html, ctx, callback
import dash_bootstrap_components as dbc
from dash.dependencies import Input, Output, State
import pandas as pd
from plotly import tools
from plotly.subplots import make_subplots
import plotly.express as px
import plotly.graph_objects as go
import numpy as np
import plotly.io as pio
import plotly.offline as pyo
import datetime
#pio.kaleido.scope.default_format = "png"
import json
import os
from pptx import Presentation
from pptx.util import Inches
import os
import logging
import pyodbc
import webbrowser
import matplotlib
matplotlib.use('TkAgg')
import matplotlib.pyplot as plt
import math
import psycopg2
#from selenium import webdriver

#matplotlib.use('agg')

#os.chdir('C:\\Users\\Dieudonne.Kantu\\source\\repos\\validproject1')
#os.chdir('C:\\Users\\Dieudonne.Kantu\\OneDrive\\Programming\\Python\\newvalidations\\validproject1')
## load data
# df = pd.read_csv('datasql1.csv')
# df1=pd.DataFrame(df)
# SQL Server connection setup
sql_server = 'ZAPL-2GC8QV3'
database = 'Validations'
username = 'localhost'
password = 'dieudonne1234'
connection_string = f'DRIVER=ODBC Driver 17 for SQL Server;SERVER={sql_server};DATABASE={database};UID={username};PWD={password}'
conn = pyodbc.connect(connection_string)
cursor = conn.cursor()

#sql_statement = "Select distinct ProjectID, CategoryName, SubCategoryName, RegionName, CountryName, AE, Share, TE from tbRun tr inner join tbCategory tcat on tcat.CategoryID = tr.CategoryID inner join tbSubCategory tscat on tscat.SubCategoryID = tr.SubCategoryID inner join tbRegion treg on treg.RegionID = tr.RegionID inner join tbCountry tcy on tcy.CountryID = tr.CountryID inner join FK_Run_Brand frb on frb.RunID = tr.RunID Group by ProjectID, CategoryName, SubCategoryName, RegionName, CountryName, AE, Share, TE Order by ProjectID, CategoryName, SubCategoryName, RegionName, CountryName, AE, Share, TE"
create_table_sql_statement = """
IF OBJECT_ID('#CorrelationTable', 'U') IS NULL
BEGIN
CREATE TABLE #CorrelationTable(
Project int,
Client int,
Category varchar(max),
SubCategory varchar(max),
Region varchar(max),
Country varchar(max),
AE decimal(18,4),
MShre decimal(18,4),
TE decimal(18,4),
SOW decimal(18,4)
)
END
"""
# Execute the SQL query
cursor.execute(create_table_sql_statement)

insert_data_sql_statement = """
INSERT into #CorrelationTable 
Select ProjectID AS Project, NULL, CategoryName AS Category, SubCategoryName AS SubCategory, RegionName AS Region, CountryName AS Country, CONVERT(decimal(18,2),AE) AS AE, CONVERT(decimal(18,2),Share) AS MShre, CONVERT(decimal(18,2),TE) AS TE, CONVERT(decimal(18,2),SOW) AS SOW
from tbRun tr left outer join tbCategory tcat on tcat.CategoryID = tr.CategoryID
  left outer join tbSubCategory tscat on tscat.SubCategoryID = tr.SubCategoryID
  left outer join tbRegion treg on treg.RegionID = tr.RegionID
  left outer join tbCountry tcy on tcy.CountryID = tr.CountryID
  left outer join FK_Run_Brand frb on frb.RunID = tr.RunID
  --Group by ProjectID, CategoryName, SubCategoryName, RegionName, CountryName, AE, Share, TE, SOW  
  --Order by ProjectID, CategoryName, SubCategoryName, RegionName, CountryName, AE, Share, TE, SOW 
"""
# Execute the SQL query
cursor.execute(insert_data_sql_statement)
# Commit the transaction
conn.commit()

insert_data_sql_statement1 = """
UPDATE #CorrelationTable
SET #CorrelationTable.Client = tbProject.ClientID
FROM #CorrelationTable
INNER JOIN tbProject ON #CorrelationTable.Project = tbProject.ProjectID
"""
# Execute the SQL query
cursor.execute(insert_data_sql_statement1)
# Commit the transaction
conn.commit()

select_data_sql_statement = """
Select * from #CorrelationTable
"""
# Execute the SQL query
cursor.execute(select_data_sql_statement)
# Fetch all rows and create a DataFrame
rows = cursor.fetchall()
merged_data = []
# Extracting specific columns from the dataset
for row in rows:
    Project = row[0]  # Assuming Project is the first column
    Client = row[1] # Assuming Client is the second column
    Category = row[2]  # Assuming Category is the third column
    SubCategory = row[3]  # Assuming SubCategory is the fourth column
    Region = row[4]
    Country = row[5]
    AE = row[6]
    MShre = row[7]
    TE = row[8]
    SOW = row[9]
    # Merge extracted columns into a single tuple
    merged_row = (Project, Client, Category, SubCategory, Region, Country, AE, TE, MShre, SOW)
    merged_data.append(merged_row)
    # Column names
    column_names = ['Project', 'Client', 'Category', 'SubCategory', 'Region', 'Country', 'AE', 'TE', 'MShre', 'SOW']
    df1 = pd.DataFrame(merged_data, columns=column_names)

select_data_sql_statement1 = """
DROP TABLE #CorrelationTable
"""
# Execute the SQL query
cursor.execute(select_data_sql_statement1)

df1['AE'] = df1['AE'].astype(float)
df1['TE'] = df1['TE'].astype(float)
df1['MShre'] = df1['MShre'].astype(float)
df1['SOW'] = df1['SOW'].astype(float)
n_obs=len(df1)
AE_Cor = df1.AE.corr(df1.MShre)
EE_Cor = df1.TE.corr(df1.MShre)
count_clients = len(np.unique(df1['Client']))
count_projects = len(np.unique(df1['Project']))

df4 = pd.DataFrame({'Correl for AE and EE' : ['Correl figures'],
                   'Cor(AE,Sh)' : [AE_Cor],
                    'Cor(EE,Sh)' : [EE_Cor],
                   'Count of brands' : [n_obs]
                   })


df2 = list(df1.Category.unique())
df2.sort()
df2a = list(df1.Category.unique())
df2a.sort()
df2.insert(0,"All categories")
df3 = list(df1.Region.unique())
df3.sort()
df3a = list(df1.Region.unique())
df3a.sort()
df3.insert(0,"All regions")
df5 = list(df1.SubCategory.unique())
df6 = list(df1.Country.unique())

#subcatlabel = []
count_countries = len(df6)
count_categories = len(df5)

df7 = pd.DataFrame({'Counts' : ['Figures'],
                   'Categories' : [count_categories],
                    'Countries' : [count_countries],
                   'Clients' : [count_clients],
                   'Projects' : [count_projects]
                   })

print("The correlation coeffient of AE and Market share: \n" , AE_Cor)
print("The correlation coeffient of EE and Market share: \n" , EE_Cor)
checklist1 = "All categories"
checklist2 = "All regions"
catlabel = "All categories"
reglabel = "All regions"


# Initialize Flask app and Dash app
server = Flask(__name__)
server.config['SECRET_KEY'] = 'your_secret_key_here'

# Create the first Dash app
app1 = dash.Dash(__name__, server=server, external_stylesheets=[dbc.themes.CYBORG], url_base_pathname='/app1/')
#app1 = dash.Dash(__name__, server=server, external_stylesheets=[dbc.themes.UNITED], url_base_pathname='/app1/')
#app1 = dash.Dash(__name__, server=server, url_base_pathname='/app1/')
app = app1.server
# Define the layout of the first Dash app
app1.layout = html.Div([
    html.H1("BVC Validations App", className="card-title"),
    dcc.Location(id='url', refresh=False),
    html.Div(id="app1-container")
])
internal_stylesheets = ['./assets/style.css']
# Create the second Dash app
#app2 = dash.Dash(__name__, server=server, suppress_callback_exceptions=True, external_stylesheets=[dbc.themes.BOOTSTRAP], url_base_pathname='/app2/')
app2 = dash.Dash(__name__, server=server, suppress_callback_exceptions=True, external_stylesheets=[dbc.themes.UNITED], url_base_pathname='/app2/')
#app2 = dash.Dash(__name__, server=server, suppress_callback_exceptions=True, external_stylesheets=internal_stylesheets, url_base_pathname='/app2/')
app = app2.server
#add validations layout here
PAGE_SIZE = 5
sidebar = html.Div(
    [
        #dbc.Row([html.P('View and tick to select')]),
        dbc.Row([html.P('Category'),
                 dcc.Checklist(
                     id="checklist1", 
                     options=[{'label': x, 'value': x} for x in df2], 
                     #options=[{'label': x, 'value': x} for x in df1.Category.unique()], 
                     value=['All categories'],
                     labelStyle={'display': 'block'})
                 ]),
        dbc.Row([html.P('Region'),
                 dcc.Checklist( 
                 id="checklist2", 
                 options=[{'label': x, 'value': x} for x in df3], 
                 value=['All regions'],
                 labelStyle={'display': 'block'})
                 ])
        #dbc.Row([
        #        dbc.Col([
        #            # represents the browser address bar and doesn't render anything
        #            dcc.Location(id='url', refresh=False),
        #            # content will be rendered in this element
        #            html.Div(id='page-content')
        #            ])

            
        #    ])
        
        ]
    )

sidebar1 = html.Div(
    [
        dbc.Row([html.P('View and press to select')]),
        dbc.Row([html.P('Category'),
                 dcc.Dropdown(
                     id="dropdown1",
                     options=[{'label': x, 'value': x} for x in df2],
                     value=['All categories'],
                     multi=True)
                 ]),
        dbc.Row([html.P('SubCategory'),
                 dcc.Dropdown(
                     id="dropdown3",
                     #options=[{'label': x, 'value': x} for x in df5],
                     value=['All subcategories'],
                     multi=True)
                 ]),
        dbc.Row([html.P('Region'),
                 dcc.Dropdown(
                     id="dropdown2",
                     options=[{'label': x, 'value': x} for x in df3],
                     value=['All regions'],
                     #value=['Africa','Latin America'],
                     multi=True)

                 ]),
        dbc.Row([html.P('Country'),
                 dcc.Dropdown(
                     id="dropdown4",
                     value=['All countries'],
                     multi=True)
            ]),
        dbc.Row([
                dbc.Col([
                    html.Br(),
                    # represents the browser address bar and doesn't render anything
                    dcc.Location(id='url', refresh=False),
                    # content will be rendered in this element
                    html.Div(id='page-content')
                    ])

            
            ])
        
        ]
    )

content = html.Div(
    [
        dbc.Row([html.P('View and download charts')]),
        dbc.Row([
            dbc.Col([dcc.Textarea(id='textarea-example',value='All categories', style={'width': '100%', 'height': 75})]),
            dbc.Col([dcc.Textarea(id='textarea-example1',value='All regions', style={'width': '100%', 'height': 75})])
            ]),
        dbc.Row([
            dbc.Col([
                #html.P('AE against Share'),
                dcc.Graph(id='scplot_share_AE')
                ]),
            dbc.Col([
                #html.P('EE(TE) against Share'),
                dcc.Graph(id='scplot_share_EE')
                ])
            ]),
        dbc.Row([
            dash_table.DataTable(
                id='datatable-paging',
                columns=[
                    #{"name": i, "id": i} for i in sorted(df4.columns)
                    {"name": i, "id": i} for i in df4.columns
               ],
               data=df4.to_dict("records"),
               page_current=0,
               page_size=PAGE_SIZE,
               page_action='native'
               )
            ]),
        dbc.Row([
            dbc.Col([
                html.Br(),
                html.Button("Download CSV", id="btn_csv"),
                dcc.Download(id="download-dataframe-csv"),
                dcc.Store(id='stored-value'),
                dcc.Store(id='stored-value3')
                ]),
            dbc.Col([
                html.Br(),
                html.Button("Download figures", id="btn_png"),
                dcc.Download(id="download_fig"),
                dcc.Store(id='stored-value1'),
                dcc.Store(id='stored-value5')
                ])
            
            ]),
        dbc.Row([
                dbc.Col([
                    html.Br(),
                    # represents the browser address bar and doesn't render anything
                    dcc.Location(id='url1', refresh=False),
                    # content will be rendered in this element
                    html.Div(id='page-content1')
                    ])

            
            ])
        ]
    
    )

content1 = html.Div(
    [
        dbc.Row([html.P('View and download tables')]),
        dbc.Row([
            dash_table.DataTable(
                id='datatable-paging1',
                columns=[
                    #{"name": i, "id": i} for i in sorted(df4.columns)
                    {"name": i, "id": i} for i in df4.columns
               ],
               data=df4.to_dict("records"),
               page_current=0,
               page_size=PAGE_SIZE,
               page_action='native'
               )
            ]),
        dbc.Row([html.P('View and download tables')]),    
        dbc.Row([
            dash_table.DataTable(
                id='datatable-paging2',
                columns=[
                    #{"name": i, "id": i} for i in sorted(df7.columns)
                    {"name": i, "id": i} for i in df7.columns
               ],
               data=df7.to_dict("records"),
               page_current=0,
               page_size=PAGE_SIZE,
               page_action='native'
               )
            ]),    
        dbc.Row([
            dbc.Col([
                html.Br(),
                html.Button("Download tables", id="btn_csv1"),
                dcc.Download(id="download-dataframe-csv1"),
                dcc.Store(id='stored-value2'), # table 1 on tab 3
                dcc.Store(id='stored-value4')  #Stores current_n_clicks. Note: n_clicks is a property attached to the button control. 
                ]),
            dbc.Col([
                html.Br(),
                html.Button("Download data", id="btn_csv2"),
                dcc.Download(id="download-dataframe-csv2"),
                dcc.Store(id='stored-value6'), # data on tab 3
                dcc.Store(id='stored-value7'), #Stores clicks_n_clicks
                ])    
            ])

    
    ])

def new_func():
    dcc.Tab([])


#end of validations layout
# Define the layout of the second Dash app
app2.layout = html.Div([
    html.H1("Validations Outputs"),
    html.P("Select categories and download files and charts."),
    dcc.Tabs(
        id="tabs-with-classes-2",
        value='tab-1',
        parent_className='custom-tabs',
        className='custom-tabs-container',
        children=[
            dcc.Tab(
                label='Tab one',
                value='tab-1',
                className='custom-tab',
                selected_className='custom-tab--selected'
            ),
            dcc.Tab(
                label='Tab two',
                value='tab-2',
                className='custom-tab',
                selected_className='custom-tab--selected',
                children=[
                    dbc.Row([dbc.Col(sidebar, width=3, className='bg-light'),
                             dbc.Col(content, width=9)],
                            style={"height": "100vh"})
                    ]
            ),
            dcc.Tab(
                label='Tab three, multiline',
                value='tab-3', className='custom-tab',
                selected_className='custom-tab--selected',
                children=[
                    dbc.Row([dbc.Col(sidebar1, width=3, className='bg-light'),
                             dbc.Col(content1, width=9)],
                            style={"height": "100vh"})
                    ]
            ),
        ]),
    html.Br(),
    html.Br(),
    #html.H2("Logout"),
   #html.Button('Logout', id='logout-button', n_clicks=0),
    dcc.Location(id='app2-url', refresh=False),
    #dcc.Link('Page 2', href='/app2/'),  # Link to /page2
    html.Div(id="app2-container"),
    html.Div(id='tabs-content-classes-2')
])

@app2.callback(Output('tabs-content-classes-2', 'children'),
              Input('tabs-with-classes-2', 'value'))
def render_content(tab):
    if tab == 'tab-1':
        return html.Div([
            #html.H3('Tab content 1'),
            html.H4('This platform gives an overview of how our equity metrics validate against real behavior. Our measures are known to outperform commonly used measures of brand equity in the following:'),
            html.Ul([  # Unordered list
            html.Li("Our measures validate well against independent sources of market share"),  # List item 1
            html.Li("Our measures validate well against surveyed measured share of wallet"),  # List item 2
            html.Li("Our measures have been validated within multiple product categories and countries"),  # List item 3
            # Empty div for vertical spacing
            html.Div(style={'height': '50px'}),  # Adjust height as needed
            html.Button('Logout', id='logout-button', n_clicks=0)
            ])
        ])
    elif tab == 'tab-2':
        return html.Div([
            #html.H3('Tab content 2')
            #html.Button('Logout', id='logout-button', n_clicks=0)
            # Empty div for vertical spacing
            html.Div(style={'height': '50px'}),  # Adjust height as needed
            html.Button('Logout', id='logout-button', n_clicks=0)
        ])
    elif tab == 'tab-3':
        return html.Div([
            #html.H3('Tab content 3')
            html.Button('Logout', id='logout-button', n_clicks=0)
        ])

# # postgreSQL Server connection setup
# # Replace these with your database and user details
# database_name = 'validations'
# user = 'postgres'
# password = 'dieudonne1234'

# # Connect to the database
# conn = psycopg2.connect(
#     database=database_name,
#     user=user,
#     password=password
# )

# # Create a cursor
# cursor = conn.cursor()

# Excel connection setup
# Read the CSV file containing the user credentials
def read_credentials():
    credentials_df = pd.read_csv('credentials.csv')  # Replace with your CSV file path
    return credentials_df.set_index('username').to_dict(orient='index')

# Login layout
# login_layout = html.Div([
#     html.H2("Please enter login details and sign into the app"),
#     html.Div(id='login-output'),
#     dcc.Input(id='username-input', type='text', placeholder='Username'),
#     dcc.Input(id='password-input', type='password', placeholder='Password'),
#     html.Button('Login', id='login-button', n_clicks=0)
# ])
login_layout = dbc.Container(
    dbc.Row(
        dbc.Col(
            dbc.Card(
                dbc.CardBody([
                    html.H2("Login", className="card-title"),
                    dcc.Input(id='EmailAddress-input',type='text',placeholder='EmailAddress',style={'margin-bottom': '10px'}),
                    dcc.Input(id='password-input',type='password',placeholder='Password',style={'margin-bottom': '10px'}),
                    html.Button('Login', id='login-button', n_clicks=0, className='btn btn-primary'),
                    html.Div(id='login-output', style={'margin-top': '10px'})
                ])
            ),
            width={'size': 6, 'offset': 3}
        ),
        style={'marginTop': '10%'}
    )
)
# Callback to handle login
@app1.callback(
    Output('login-output', 'children', allow_duplicate=True),
    Output('url', 'pathname', allow_duplicate=True),
    Input('login-button', 'n_clicks'),
    State('EmailAddress-input', 'value'),
    State('password-input', 'value'),
    prevent_initial_call=True
)
def handle_login(n_clicks, EmailAddress, password):

    if n_clicks > 0:
        # user_credentials = read_credentials()
        # if username in user_credentials and user_credentials[username]['password'] == password:
        # Select user credentials based on the EmailAddress and password
        cursor.execute("SELECT * FROM tbUser WHERE EmailAddress = ? AND password = ?", (EmailAddress, password))
        user = cursor.fetchone()
        if user:
            session['logged_in'] = True
            return "Logged in successfully!", "/app1/"
        else:
            return "Invalid credentials", "/"
    return "", ""


## Callback to handle not logged in
@app2.callback(
    Output('app2-container', 'children', allow_duplicate=True),
    #Output('app2-url', 'pathname', allow_duplicate=True),
    Input('app2-url', 'pathname'),
    prevent_initial_call=True
)
def handle_notloggedin(pathname):
    if 'logged_in' not in session and pathname == '/app2/':
        return  dcc.Location(id='url', pathname='/app1/')
    elif 'logged_in' in session and pathname == '/app2/':
        #return app2.layout, "/app2/"
        return ""


# Callback to display login page if not logged in or redirect to second app if logged in
@app1.callback(Output('app1-container', 'children'), 
               Input('url', 'pathname')
)
def display_page(pathname):
    if 'logged_in' not in session and pathname == '/app1/':
        return login_layout
    elif 'logged_in' in session and pathname == '/app1/':
        return dcc.Location(id='app2-url', pathname='/app2/')
    return login_layout


# Callback to handle logout
@app2.callback(
    Output('app2-container', 'children'),
    Output('logout-button', 'n_clicks', allow_duplicate=True),
    Input('logout-button', 'n_clicks'),
    #Input('app2-url', 'pathname'),
    prevent_initial_call=True
)
def handle_logout(n_clicks):
    if n_clicks > 0:
        session.pop('logged_in', None)
        return dcc.Location(id='url', pathname='/app1/'), 0
    return dcc.Location(id='app2-url', pathname='/app2/'), 0


@app2.callback(
    Output('textarea-example', 'value'),
    Output('textarea-example1', 'value'),
    Output(component_id='scplot_share_AE', component_property='figure'),
    Output(component_id='scplot_share_EE', component_property='figure'),
    Output('datatable-paging', 'data'),
    Output('stored-value', 'data'),
    Output('stored-value1', 'data'),
    Input(component_id='checklist1', component_property='value'),
    Input(component_id='checklist2', component_property='value'),
    Input('textarea-example', 'value'),
    Input('textarea-example1', 'value'),
    Input('datatable-paging', "page_current"),
    Input('datatable-paging', "page_size")
    )
def update_graph(checklist1,checklist2,value,value1,page_current,page_size):
# You can update Textarea by typing the values in its window.  Textarea has both an input and output part. 
    #pathname = ""
    value = '\n'.join(checklist1)
    value1 = '\n'.join(map(str,checklist2))
  # If statements for categories and regions
    if 'All categories' in checklist1:
        checklist1 = df2a
    # elif checklist1 is None:
    #     checklist1 = df2a
    elif len(checklist1) == 0:  
        checklist1 = df2a
    else:
        checklist1 = checklist1

    if 'All regions' in checklist2:
        checklist2 = df3a
    # elif checklist2 is None:
    #     checklist2 = df3a
    elif len(checklist2) == 0:
        checklist2 = df3a
    else:
        checklist2 = checklist2

    df_filter = df1.loc[(df1['Category'].isin(checklist1)) & (df1['Region'].isin(checklist2))]

    n_obs=len(df_filter)
    AE_Cor = df_filter.AE.corr(df_filter.MShre)
    EE_Cor = df_filter.TE.corr(df_filter.MShre)
    mydf = pd.DataFrame({'Correl for AE and EE' : ['Correl figures'],
                   'Cor(AE,Sh)' : [round(AE_Cor,2)],
                    'Cor(EE,Sh)' : [round(EE_Cor,2)],
                   'Count of brands' : [n_obs]
                   })

    df4 = mydf.iloc[page_current*page_size:(page_current+ 1)*page_size].to_dict('records')

    if(len(df_filter)<=0):
        fig1 = go.Figure().add_annotation(x=2, y=2,text="No Data to Display",font=dict(family="sans serif",size=25,color="crimson"),showarrow=False,yshift=10)
        fig1.update_layout(title="AE against market share",xaxis_title="Attitudinal Equity", yaxis_title="Market share")
        fig2 = go.Figure().add_annotation(x=2, y=2,text="No Data to Display",font=dict(family="sans serif",size=25,color="crimson"),showarrow=False,yshift=10)
        fig2.update_layout(title="EE(TE) against market share",xaxis_title="Effective Equity", yaxis_title="Market share")
    else:
        fig1 = px.scatter(df_filter, x="AE", y="MShre", trendline="lowess", labels={"AE":"Attitudinal Equity", "MShre":"Market share"}, title='AE against market share') #Error: Invalid value
        fig2 = px.scatter(df_filter, x="TE", y="MShre", trendline="ols", labels={"TE":"Effective Equity", "MShre":"Market share"}, title='EE(TE) against market share') #Error: Invalid value

    return value, value1, fig1, fig2, df4, mydf.to_json(date_format='iso', orient='split'), df_filter.to_json(date_format='iso', orient='split')

#return '\n{}'.format(value), '\n{}'.format(value1)

@app2.callback(
    Output("download-dataframe-csv", "data"),
    Output("btn_csv", "n_clicks", allow_duplicate=True),
    Output('stored-value3', 'data', allow_duplicate=True),
    Input("btn_csv", "n_clicks"),
    Input('stored-value', 'data'),
    Input('stored-value3', 'data'),
    prevent_initial_call=True
)
def execute_command(n_clicks,jsonified_cleaned_data,current_n_clicks):
    if n_clicks is None:
        n_clicks = 0
        current_n_clicks = 0
        return "", n_clicks, current_n_clicks

    if current_n_clicks is None:
        #n_clicks = 0
        current_n_clicks = 0
        return "", n_clicks, current_n_clicks
# current_n_clicks is created to prevent a download of "Validations_table.csv" each time there is a new selection (this was the observed behavior). The download is only activated when n_clicks > current_n_clicks. 
    if n_clicks is not None and current_n_clicks is not None:
        if isinstance(current_n_clicks, str):
            current_n_clicks = 0          
            
        if n_clicks > current_n_clicks:
            # Put your command to execute here
            dff = pd.read_json(jsonified_cleaned_data, orient='split')
            # Update current_n_clicks to match the new n_clicks value
            current_n_clicks = n_clicks
            return dcc.send_data_frame(dff.to_csv, "Validations_table.csv"), n_clicks, current_n_clicks
        else:
            return "", n_clicks, current_n_clicks

@app2.callback(
    Output('page-content', 'children'),
    [Input('url', 'pathname')])
def display_page(relative_pathname):
    return html.Div([
        #html.H3(f'You are on page {relative_pathname}'),
        html.A(html.Button('Refresh Page'),href=relative_pathname),
    ])

@app2.callback(
    Output('page-content1', 'children'),
    [Input('url1', 'pathname')])
def display_page(relative_pathname):
    return html.Div([
        #html.H3(f'You are on page {relative_pathname}'),
        html.A(html.Button('Refresh Page'),href=relative_pathname),
    ])



@app2.callback(
    Output("download_fig", "data"),
    Output("btn_png", "n_clicks", allow_duplicate=True),
    Output('stored-value5', 'data', allow_duplicate=True),
    Input("btn_png","n_clicks"),
    Input('stored-value1', 'data'),
    Input('stored-value5', 'data'),
    prevent_initial_call=True
)

def func(n_clicks,jsonified_cleaned_data1,current_n_clicks):
    if n_clicks is None:
        n_clicks = 0
        current_n_clicks = 0
        return "", n_clicks, current_n_clicks
        
    if current_n_clicks is None:
        #n_clicks = 0
        current_n_clicks = 0
        return "", n_clicks, current_n_clicks

    #if n_clicks is not None and current_n_clicks is not None:
    if n_clicks is not None and current_n_clicks is not None:
        if isinstance(current_n_clicks, str):
            current_n_clicks = 0    
        if n_clicks > current_n_clicks:
            dff = pd.read_json(jsonified_cleaned_data1, orient='split')
            df_filter = dff
            now_date = datetime.datetime.now()
            time_today = now_date.strftime("%d-%m-%Y_%H-%M-%S-%f")
            #fig1_path = Path("./assets/fig1_" + time_today + ".png")
            fig1_path = os.path.normpath("./assets/fig1_" + time_today + ".png")
            fig2_path = os.path.normpath("./assets/fig2_" + time_today + ".png")
            if(len(df_filter)<=0):
                fig1 = go.Figure().add_annotation(x=2, y=2,text="No Data to Display",font=dict(family="sans serif",size=25,color="crimson"),showarrow=False,yshift=10)
                fig2 = go.Figure().add_annotation(x=2, y=2,text="No Data to Display",font=dict(family="sans serif",size=25,color="crimson"),showarrow=False,yshift=10)
            else:
                fig1 = px.scatter(df_filter, x=df_filter.AE, y=df_filter.MShre, trendline="lowess") #Error: Invalid value
                fig2 = px.scatter(df_filter, x=df_filter.TE, y=df_filter.MShre, trendline="ols") #Error: Invalid value
            # Save figures
            ## Convert the HTML to an image and save it
            #fig1.to_image(format='png', width=800, height=600, scale=2).write_image(fig1_path)
            #pyo.plot(fig2, filename='temp_plot2.html', auto_open=False)
            #fig2.to_image(format='png', width=800, height=600, scale=2).write_image(fig2_path)
            pio.write_image(fig1, fig1_path, format='png')
            pio.write_image(fig2, fig2_path, format='png')
            # create an Object
            ppt = Presentation()
            first_slide = ppt.slides.add_slide(ppt.slide_layouts[0])
            # title (included date)
            title = "Dowload validations charts"
            # set the title on first slide
            first_slide.shapes[0].text_frame.paragraphs[0].text = title
            # slide 2 - set the image
            #img = 'C:/Users/DKantu01/PythonApps/BasicProject/BasicProject/assets/fig1.png'
            second_slide = ppt.slide_layouts[1]
            slide2 = ppt.slides.add_slide(second_slide)
            # play with the image attributes if you are not OK with the height and width
            pic = slide2.shapes.add_picture(fig1_path, left= Inches(2),top = Inches(1),height = Inches(5))
            # slide 3 - set the image
            third_slide = ppt.slide_layouts[1]
            slide3 = ppt.slides.add_slide(third_slide)
            # play with the image attributes if you are not OK with the height and width
            pic = slide3.shapes.add_picture(fig2_path, left= Inches(2),top = Inches(1),height = Inches(5))
            # save the powerpoint presentation
            ppt_path = os.path.normpath("./assets/Validations_charts_" + time_today + ".pptx")
            ppt.save(ppt_path)
            # Close the figures if needed
            #plt.close(fig1_path)
            #plt.close(fig2_path)
            #return dcc.send_file("./assets/fig1.png")
            #return dcc.send_file(fig1_path)
            current_n_clicks = n_clicks
            return dcc.send_file(ppt_path), n_clicks, current_n_clicks



# Create a callback from Category & Subcategory dropdowns and Region and Country Dropdowns
@app2.callback(
    #Output('stored-value2', 'data'),
    Output('dropdown3', 'options'),
    Output('dropdown3', 'value'),
    Output('dropdown4', 'options'),
    Output('dropdown4', 'value'),
    #Output('datatable-paging1', 'data'),
    Input('dropdown1', 'value'),
    Input('dropdown2', 'value')
    #Input('dropdown3', 'value'),
    #Input('datatable-paging1', "page_current"),
    #Input('datatable-paging1', "page_size")
    )

def update_subcatcountry(catlabel,reglabel):
    if ('All categories' in catlabel and 'All regions' in reglabel):
        subcatlabel = df5
        catlabel = df2a
        countrylabel = df6
        reglabel = df3a
        relevant_subcat_options = []
        relevant_subcat_options.insert(0,"All subcategories")
        relevant_country_options = []
        relevant_country_options.insert(0,"All countries")
    # Create and return formatted relevant options with the same label and value
        formatted_relevant_subcat_options = [{'label':x, 'value':x} for x in relevant_subcat_options if x is not None]
        formatted_relevant_country_options = [{'label':x, 'value':x} for x in relevant_country_options if x is not None]
    elif ('All categories' not in catlabel and 'All regions' in reglabel):
        countrylabel = df6
        reglabel = df3a
        relevant_country_options = []
        relevant_country_options.insert(0,"All countries")
        df_filter = df1.loc[(df1['Category'].isin(catlabel)) & (df1['Region'].isin(reglabel))]
        subcatlabel = list(df_filter.loc[(df_filter['Category'].isin(catlabel))]['SubCategory'].unique())
        subcatlabel.sort()
        formatted_relevant_country_options = [{'label':x, 'value':x} for x in relevant_country_options]
        formatted_relevant_subcat_options = [{'label':x, 'value':x} for x in subcatlabel]
    elif ('All categories' in catlabel and 'All regions' not in reglabel):
        subcatlabel = df5
        catlabel = df2a
        relevant_subcat_options = []
        relevant_subcat_options.insert(0,"All subcategories")
        df_filter = df1.loc[(df1['Category'].isin(catlabel)) & (df1['Region'].isin(reglabel))]
        countrylabel = list(df_filter.loc[(df_filter['Region'].isin(reglabel))]['Country'].unique())
        countrylabel.sort()
        formatted_relevant_subcat_options = [{'label':x, 'value':x} for x in relevant_subcat_options]
        formatted_relevant_country_options = [{'label':x, 'value':x} for x in countrylabel]
    elif ('All categories' not in catlabel and 'All regions' not in reglabel):
        df_filter = df1.loc[(df1['Category'].isin(catlabel)) & (df1['Region'].isin(reglabel))]
        countrylabel = list(df_filter.loc[(df_filter['Region'].isin(reglabel))]['Country'].unique())
        subcatlabel = list(df_filter.loc[(df_filter['Category'].isin(catlabel))]['SubCategory'].unique())
        countrylabel.sort()
        subcatlabel.sort()
        formatted_relevant_country_options = [{'label':x, 'value':x} for x in countrylabel]
        formatted_relevant_subcat_options = [{'label':x, 'value':x} for x in subcatlabel]
    else:
        subcatlabel = list(df1.loc[(df1['Category'].isin(catlabel))]['SubCategory'].unique())
        subcatlabel.sort()
        formatted_relevant_subcat_options = [{'label':x, 'value':x} for x in subcatlabel]
        countrylabel = list(df1.loc[(df1['Region'].isin(reglabel))]['Country'].unique())
        countrylabel.sort()
        formatted_relevant_country_options = [{'label':x, 'value':x} for x in countrylabel]
    
    return formatted_relevant_subcat_options, formatted_relevant_subcat_options, formatted_relevant_country_options, formatted_relevant_country_options
    

@app2.callback(
    Output('datatable-paging1', 'data'),
    Output('datatable-paging2', 'data'),
    Output('stored-value2', 'data'),
    Output('stored-value5', 'data'), #stores second table on tab 3 with category and client counts
    Output('stored-value6', 'data'), #stores data on tab 3 that is used to generate tables
    Input('dropdown1', 'value'),
    Input('dropdown2', 'value'),
    Input('dropdown3', 'value'),
    Input('dropdown4', 'value'),
    Input('datatable-paging1', "page_current"),
    Input('datatable-paging1', "page_size"))

    #def update_subcategory(catlabel,reglabel, subcatlabel, page_current,page_size):

def update_table_tab3(catlabel,reglabel,subcatlabel,countrylabel,page_current,page_size):
    if ('All categories' in catlabel and 'All regions' in reglabel):
        subcatlabel = df5
        catlabel = df2a
        countrylabel = df6
        reglabel = df3a
    elif ('All categories' not in catlabel and 'All regions' in reglabel):
        catlabel = catlabel
        countrylabel = df6
        reglabel = df3a
        df_filter = df1.loc[(df1['Category'].isin(catlabel)) & (df1['Region'].isin(reglabel)) & (df1['Country'].isin(countrylabel))]
        listsubcat = list(df_filter.loc[(df_filter['Category'].isin(catlabel))]['SubCategory'].unique())
        if (len(subcatlabel)== len(listsubcat)):
          subcatlabel = df5
        else:
          subcatlabel = subcatlabel

    elif ('All categories' in catlabel and 'All regions' not in reglabel):
        subcatlabel = df5
        catlabel = df2a
        reglabel = reglabel
        df_filter = df1.loc[(df1['Category'].isin(catlabel)) & (df1['SubCategory'].isin(subcatlabel)) & (df1['Region'].isin(reglabel)) ]
        listcountry = list(df_filter.loc[(df_filter['Region'].isin(reglabel))]['Country'].unique())
        if (len(countrylabel) == len(listcountry)):
          countrylabel = df6
        else:
          countrylabel = countrylabel
        
    elif ('All categories' not in catlabel and 'All regions' not in reglabel):
        catlabel = catlabel
        reglabel = reglabel
        df_filter = df1.loc[(df1['Category'].isin(catlabel)) & (df1['Region'].isin(reglabel))]
        listcountry = list(df_filter.loc[(df_filter['Region'].isin(reglabel))]['Country'].unique())
        listsubcat = list(df_filter.loc[(df_filter['Category'].isin(catlabel))]['SubCategory'].unique())

        if (len(subcatlabel) == len(listsubcat)):
          subcatlabel = df5
        else:
          subcatlabel = subcatlabel
        
        if (len(countrylabel) == len(listcountry)):
          countrylabel = df6
        else:
          countrylabel = countrylabel
   
    else:
        subcatlabel = subcatlabel
        catlabel = catlabel
        countrylabel = countrylabel
        reglabel = reglabel
        df_filter = df1.loc[(df1['Category'].isin(catlabel)) & (df1['SubCategory'].isin(subcatlabel)) & (df1['Region'].isin(reglabel)) & (df1['Country'].isin(countrylabel))]

        
    df_filter = df1.loc[(df1['Category'].isin(catlabel)) & (df1['SubCategory'].isin(subcatlabel)) & (df1['Region'].isin(reglabel)) & (df1['Country'].isin(countrylabel))]

    n_obs=len(df_filter)
    AE_Cor = df_filter.AE.corr(df_filter.MShre)
    EE_Cor = df_filter.TE.corr(df_filter.MShre)
    client = df_filter['Client']
    project = df_filter['Project']
    country = df_filter['Country']
    subCategory = df_filter['SubCategory']

    client_cleaned = [item for item in client if item is not None]
    project_cleaned = [item for item in project if item is not None]
    country_cleaned = [item for item in country if item is not None]
    subCategory_cleaned = [item for item in subCategory if item is not None]
   
    count_clients = len(np.unique(client_cleaned))
    count_projects = len(np.unique(project_cleaned))
    count_countries = len(np.unique(country_cleaned))
    count_categories = len(np.unique(subCategory_cleaned))
   
    # count_clients = len(np.unique(df_filter['Client']))
    # count_projects = len(np.unique(df_filter['Project']))
    # count_countries = len(np.unique(df_filter['Country']))
    # count_categories = len(np.unique(df_filter['SubCategory']))


    mydf = pd.DataFrame({'Correl for AE and EE' : ['Correl figures'],
                   'Cor(AE,Sh)' : [round(AE_Cor,2)],
                    'Cor(EE,Sh)' : [round(EE_Cor,2)],
                   'Count of brands' : [n_obs]
                   })
    
    mydf1 = pd.DataFrame({'Counts' : ['Count figures'],
                   'Categories' : [count_categories],
                    'Countries' : [count_countries],
                   'Clients' : [count_clients],
                   'Projects' : [count_projects]
                   })

    df4 = mydf.iloc[page_current*page_size:(page_current+ 1)*page_size].to_dict('records')
    df7 = mydf1.iloc[page_current*page_size:(page_current+ 1)*page_size].to_dict('records')

    return df4, df7, mydf.to_json(date_format='iso', orient='split'), mydf1.to_json(date_format='iso', orient='split'), df_filter.to_json(date_format='iso', orient='split')


@app2.callback(
    Output("download-dataframe-csv1", "data"),
    Output("btn_csv1", "n_clicks", allow_duplicate=True),
    Output('stored-value4', 'data', allow_duplicate=True),
    Input("btn_csv1", "n_clicks"),
    Input('stored-value2', 'data'),# table 1 on tab 3
    Input('stored-value5', 'data'),# table 2 on tab 3
    Input('stored-value4', 'data'),# Current_n_clicks
    prevent_initial_call=True
)
def execute_command(n_clicks,jsonified_cleaned_data,jsonified_cleaned_data1, current_n_clicks):
    if n_clicks is None:
        n_clicks = 0
        current_n_clicks = 0
        return "", n_clicks, current_n_clicks

    if current_n_clicks is None:
        #n_clicks = 0
        current_n_clicks = 0
        return "", n_clicks, current_n_clicks

    if n_clicks is not None and current_n_clicks is not None:
        if isinstance(current_n_clicks, str):
            current_n_clicks = 0    
        if n_clicks > current_n_clicks:
            # Put your command to execute here
            dff = pd.read_json(jsonified_cleaned_data, orient='split') # read table 1 on tab 3 from stored-value2
            dff1 = pd.read_json(jsonified_cleaned_data1, orient='split') # read table 2 on tab 3 from stored-value5
            # Concatenate the DataFrames vertically
            dff = pd.concat([dff, dff1], ignore_index=True)
            # Update current_n_clicks to match the new n_clicks value
            current_n_clicks = n_clicks
            return dcc.send_data_frame(dff.to_csv, "Validations_table.csv"), n_clicks, current_n_clicks
        else:
            return "", n_clicks, current_n_clicks
        
@app2.callback(
    Output("download-dataframe-csv2", "data"),
    Output("btn_csv2", "n_clicks", allow_duplicate=True),
    Output('stored-value7', 'data', allow_duplicate=True),
    Input("btn_csv2", "n_clicks"),
    Input('stored-value6', 'data'),# data on tab 3
    Input('stored-value7', 'data'),# Current_n_clicks
    prevent_initial_call=True
)
def execute_command(n_clicks,jsonified_cleaned_data, current_n_clicks):
    if n_clicks is None:
        n_clicks = 0
        current_n_clicks = 0
        return "", n_clicks, current_n_clicks

    if current_n_clicks is None:
        #n_clicks = 0
        current_n_clicks = 0
        return "", n_clicks, current_n_clicks

    if n_clicks is not None and current_n_clicks is not None:
        if isinstance(current_n_clicks, str):
            current_n_clicks = 0    
        if n_clicks > current_n_clicks:
            # Put your command to execute here
            dff = pd.read_json(jsonified_cleaned_data, orient='split')
            # Update current_n_clicks to match the new n_clicks value
            current_n_clicks = n_clicks
            return dcc.send_data_frame(dff.to_csv, "Data_table.csv"), n_clicks, current_n_clicks
        else:
            return "", n_clicks, current_n_clicks

app.debug = True
app.logger.setLevel(logging.DEBUG)

# if __name__ == '__main__':
#    app.run()

if __name__ == '__main__':  
    #app.run(debug=True)  
    import os 
    HOST = os.environ.get('SERVER_HOST', 'localhost') 
    try: 
        PORT = int(os.environ.get('SERVER_PORT', '8051')) 
    except ValueError: 
        PORT = 8051 
    server.run(HOST, PORT) 
    server.run(debug=True)

# if __name__ == '__main__':
#     port = int(os.environ.get("PORT", 8000))
#     app.run(host='0.0.0.0', port=port)

# if __name__ == '__main__':
#     #app.run(port=60000)
#     app.run_server(port=8000)
